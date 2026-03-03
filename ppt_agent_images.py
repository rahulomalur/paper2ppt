import os
import tarfile
import zipfile
import streamlit as st
import io
import contextlib
from langchain.tools import tool 
from langchain_text_splitters.latex import LatexTextSplitter
from langchain_groq import ChatGroq
from langchain.agents import initialize_agent, AgentType
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
from pydantic import BaseModel, Field, validator
from typing import List, Optional
from pptx.dml.color import RGBColor

load_dotenv()

# --- INITIALIZE SESSION STATE ---
if "step" not in st.session_state:
    st.session_state.step = "upload"
if "project_dir" not in st.session_state:
    st.session_state.project_dir = None
if "outline" not in st.session_state:
    st.session_state.outline = None
if "ppt_path" not in st.session_state:
    st.session_state.ppt_path = None
if "agent_thoughts" not in st.session_state:
    st.session_state.agent_thoughts = ""

# --- SIDEBAR LOGS ---
with st.sidebar:
    st.header(" Agent Reasoning")
    if st.session_state.agent_thoughts:
        st.text_area("Internal Monologue", st.session_state.agent_thoughts, height=500)

# --- HELPERS ---
def extract_project_files(uploaded_file):
    """Unzips .zip or .tar.gz and targets ONLY main.tex."""
    project_path = "extracted_project"
    if not os.path.exists(project_path):
        os.makedirs(project_path)
    
    file_name = uploaded_file.name
    
    # Handle .zip files
    if file_name.endswith(".zip"):
        with zipfile.ZipFile(uploaded_file, 'r') as zip_ref:
            zip_ref.extractall(project_path)
    # Handle .tar.gz files
    elif file_name.endswith(".tar.gz") or file_name.endswith(".tgz"):
        with tarfile.open(fileobj=uploaded_file, mode="r:gz") as tar:
            tar.extractall(path=project_path)
    
    # Search specifically for main.tex
    for root, dirs, files in os.walk(project_path):
        if "main.tex" in files:
            file_path = os.path.join(root, "main.tex")
            with open(file_path, 'r', errors='ignore') as f:
                return f.read(), project_path
    
    return None, project_path

# --- TOOLS ---
class Slide(BaseModel):
    title: str = Field(description="Title of the slide")
    points: List[str] = Field(description="List of bullet points.")
    image_path: Optional[str] = Field(None, description="Filename of a figure (e.g., 'chart.png')")

    @validator('points', pre=True)
    def ensure_list(cls, v):
        if isinstance(v, str): return [v]
        return v

class PPTInput(BaseModel):
    slides: List[Slide] = Field(description="List of slide objects")

@tool(args_schema=PPTInput)
def ppt_create(slides: List[Slide]) -> str:
    """Creates PPTX with a hard rule: slides are either text-only or image-only."""
    prs = Presentation()
    project_dir = st.session_state.get("project_dir", ".")
    
    for s_data in slides:
        # If points are empty, use a 'Title Only' layout (Layout 5)
        is_visual_only = len(s_data.points) == 0
        layout_idx = 5 if is_visual_only else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide.shapes.title.text = s_data.title
        
        # 1. Text Logic: Only add if points exist
        if not is_visual_only:
            tf = slide.shapes.placeholders[1].text_frame
            for pt in s_data.points:
                tf.add_paragraph().text = str(pt)
        
        # 2. Image Logic: Only add if an image_path is provided
        if s_data.image_path:
            clean_name = os.path.basename(s_data.image_path)
            found_path = None
            for r, _, fs in os.walk(project_dir):
                if clean_name in fs:
                    found_path = os.path.join(r, clean_name)
                    break
            
            if found_path:
                # Center and expand the image since there is no text
                # Left, Top, Width (centered on 10-inch wide slide)
                prs.slides[prs.slides.index(slide)].shapes.add_picture(
                    found_path, Inches(1), Inches(1.5), width=Inches(8)
                )
    
    path = os.path.abspath("output_presentation.pptx")
    prs.save(path)

@tool
def beautify_presentation(file_path: str, theme: str) -> str:
    """Applies theme colors and fonts."""
    try:
        prs = Presentation(file_path)
        configs = {
            "Corporate": {"bg": (0, 32, 96), "text": (255, 255, 255), "font": "Arial"},
            "Modern": {"bg": (240, 240, 240), "text": (40, 40, 40), "font": "Segoe UI"},
            "Research": {"bg": (255, 255, 255), "text": (0, 0, 0), "font": "Times New Roman"}
        }
        style = configs.get(theme, configs["Modern"])
        for slide in prs.slides:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*style["bg"])
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.name = style["font"]
                            run.font.color.rgb = RGBColor(*style["text"])
        prs.save(file_path)
        return f"Applied {theme}."
    except Exception as e: return f"Error: {str(e)}"

# --- AGENT LOGIC ---
def run_summarizer(full_latex_content, model_name):
    """Processes main.tex with a progress bar and research filter."""
    llm = ChatGroq(model=model_name, api_key=os.getenv("GROQ_API_KEY_2"), temperature=0)
    splitter = LatexTextSplitter(chunk_size=4000, chunk_overlap=200)
    chunks = [doc.page_content for doc in splitter.create_documents([full_latex_content])]
    
    summaries = []
    progress_text = st.empty()
    progress_bar = st.progress(0)
    total_chunks = len(chunks)

    for i, chunk in enumerate(chunks):
        current_progress = (i + 1) / total_chunks
        progress_text.text(f"Summarizing main.tex segment {i+1} of {total_chunks}...")
        progress_bar.progress(current_progress)

        res = llm.invoke(f"""
    Extract the ACTUAL research data from this LaTeX segment.
    
    RULES:
    - DO NOT say "This section describes..." or "The slide should show...".
    - DO NOT give instructions.
    - Include Image Paths if figures are mentioned with the correct path (e.g., 'figures/chart.png').
    
    
    Segment: {chunk}
""")
        if len(res.content) > 50:
            summaries.append(res.content)
    
    progress_text.empty()
    progress_bar.empty()

    master_outline = llm.invoke("Combine these summaries into a 10-15 slide outline. Include 'image_path' if a figure was found:\n" + "\n".join(summaries))
    st.session_state.agent_thoughts += f"\n--- SUMMARIZER ---\nParsed main.tex into outline."
    return master_outline.content

def run_creator(outline, model_name):
    llm = ChatGroq(model=model_name, api_key=os.getenv("GROQ_API_KEY_2"), temperature=0)
    agent = initialize_agent([ppt_create], llm, agent=AgentType.STRUCTURED_CHAT_ZERO_SHOT_REACT_DESCRIPTION, verbose=True)
    
    f = io.StringIO()
    with contextlib.redirect_stdout(f):
        prompt = f"""
Using the outline provided, call 'ppt_create_with_images' to build the PPT.

STRICT LAYOUT RULES:
1. For every section in the outline that contains an 'image_path', you MUST create TWO slides:
   - SLIDE A (Content): Title + Bullet points explaining the data. NO image here.
   - SLIDE B (Visual): The same Title + ONE image_path. NO points here (points should be an empty list []).
2. If an image is found, SLIDE B must be completely empty of text except for the Title.
3. Every single image found in the outline must get its own dedicated 'Visual' slide.

SCHEMA EXAMPLE FOR A FIGURE:
{{
  "slides": [
    {{
      "title": "TOPSIS Results",
      "points": ["Actual data point 1", "Actual data point 2"]
    }},
    {{
      "title": "TOPSIS Results (Visual)",
      "points": [],
      "image_path": "topsis_results.png"
    }}
  ]
}}

Outline: {outline}
"""
        agent.run(prompt)
    st.session_state.agent_thoughts += f"\n--- CREATOR ---\n{f.getvalue()}"

def run_designer(ppt_path, model_name):
    llm = ChatGroq(model='llama-3.3-70b-versatile', api_key=os.getenv("GROQ_API_KEY_2"), temperature=0)
    agent = initialize_agent([beautify_presentation], llm, agent=AgentType.STRUCTURED_CHAT_ZERO_SHOT_REACT_DESCRIPTION, verbose=True)
    f = io.StringIO()
    with contextlib.redirect_stdout(f):
        agent.run(f"Pick the best theme (Corporate, Modern, Research) for {ppt_path} and apply it.")
    st.session_state.agent_thoughts += f"\n--- DESIGNER ---\n{f.getvalue()}"

# --- UI ---
st.title("Research Paper to PPT")
model_choice = st.selectbox("LLM Selection", ("llama-3.3-70b-versatile", "openai/gpt-oss-20b","openai/gpt-oss-120b", "qwen/qwen3-32b"))
uploaded_file = st.file_uploader("Upload .zip or .tar.gz Archive", type=["zip", "gz"])

if uploaded_file and st.session_state.step == "upload":
    main_content, p_dir = extract_project_files(uploaded_file)
    if main_content:
        st.session_state.project_dir = p_dir
        if st.button("Analyze main.tex"):
            st.session_state.outline = run_summarizer(main_content, model_choice)
            st.session_state.step = "approve"
            st.rerun()
    else:
        st.error("No 'main.tex' found in the project.")

if st.session_state.step == "approve":
    st.subheader("Review Research Outline")
    with st.expander("View Outline", expanded=True):
        st.markdown(st.session_state.outline)
    col1, col2 = st.columns(2)
    with col1:
        if st.button("Approve & Create"):
            with st.spinner("Generating PowerPoint..."):
                run_creator(st.session_state.outline, model_choice)
                st.session_state.ppt_path = os.path.abspath("output_presentation.pptx")
                st.session_state.step = "design"
                st.rerun()
    with col2:
        if st.button("Restart"):
            st.session_state.clear()
            st.rerun()

if st.session_state.step == "design":
    st.info("Slides created. Ready for theme application.")
    if st.button("✨ Apply Theme"):
        with st.spinner("Beautifying..."):
            run_designer(st.session_state.ppt_path, model_choice)
            st.session_state.step = "final"
            st.rerun()

if st.session_state.step == "final":
    st.success("Complete!")
    with open(st.session_state.ppt_path, "rb") as f:
        st.download_button("📥 Download PPT", f, file_name="presentation.pptx")
    if st.button("New Paper"):
        st.session_state.clear()
        st.rerun()