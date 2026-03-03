import os
import streamlit as st
import io
import contextlib
from langchain.tools import tool 
from langchain_text_splitters.latex import LatexTextSplitter
from langchain_groq import ChatGroq
from langchain.agents import initialize_agent, AgentType
from pptx import Presentation
from dotenv import load_dotenv
from pydantic import BaseModel, Field, validator
from typing import List
from pptx.util import Pt
from pptx.dml.color import RGBColor

load_dotenv()

# Support both local .env and Streamlit Cloud secrets
GROQ_API_KEY = os.getenv("GROQ_API_KEY_1") or st.secrets.get("GROQ_API_KEY_1", "")

# --- INITIALIZE SESSION STATE ---
if "step" not in st.session_state:
    st.session_state.step = "upload"
if "outline" not in st.session_state:
    st.session_state.outline = None
if "ppt_path" not in st.session_state:
    st.session_state.ppt_path = None
if "agent_thoughts" not in st.session_state:
    st.session_state.agent_thoughts = ""

# --- SIDEBAR LOGS ---
with st.sidebar:
    st.header("🧠 Agent Reasoning")
    if st.session_state.agent_thoughts:
        st.text_area("Internal Monologue", st.session_state.agent_thoughts, height=500)

# --- TOOLS ---
@tool
def beautify_presentation(file_path: str, theme: str) -> str:
    """Applies Corporate, Modern, or Research theme to a PPTX file."""
    try:
        prs = Presentation(file_path)
        configs = {
            "Corporate": {"bg": RGBColor(0, 32, 96), "title_color": RGBColor(255, 255, 255), "body_color": RGBColor(200, 200, 200), "font": "Arial"},
            "Modern": {"bg": RGBColor(240, 240, 240), "title_color": RGBColor(40, 40, 40), "body_color": RGBColor(60, 60, 60), "font": "Segoe UI"},
            "Research": {"bg": RGBColor(255, 255, 255), "title_color": RGBColor(0, 0, 0), "body_color": RGBColor(30, 30, 30), "font": "Times New Roman"}
        }
        style = configs.get(theme, configs["Modern"])
        for slide in prs.slides:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = style["bg"]
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = style["font"]
                            run.font.color.rgb = style["title_color"] if shape == slide.shapes.title else style["body_color"]
        prs.save(file_path)
        return f"Successfully applied {theme} theme."
    except Exception as e:
        return f"Error: {str(e)}"

class Slide(BaseModel):
    title: str = Field(description="Title of the slide")
    points: List[str] = Field(description="List of bullet points as strings.")

    @validator('points', pre=True)
    def ensure_list(cls, v):
        if isinstance(v, str): return [v]
        return v

class PPTInput(BaseModel):
    slides: List[Slide] = Field(description="List of slide objects")

@tool
def latex_parse(latex_file: str) -> list:
    """Parses .tex file into manageable chunks."""
    try:
        with open(latex_file, 'r') as file:
            content = file.read()
        splitter = LatexTextSplitter(chunk_size=4000, chunk_overlap=200)
        return [doc.page_content for doc in splitter.create_documents([content])]
    except Exception as e:
        return [f"Error reading file: {str(e)}"]

@tool(args_schema=PPTInput)
def ppt_create(slides: List[Slide]) -> str:
    """Creates PPTX from structured slides."""
    prs = Presentation()
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.title
        tf = slide.shapes.placeholders[1].text_frame
        for point in slide_data.points:
            p = tf.add_paragraph()
            p.text = str(point)
    path = os.path.abspath("output_presentation.pptx")
    prs.save(path)
    return path

# --- CORE FUNCTIONS ---
def run_summarizer(latex_file, model_name):
    """Agent 1: Reads long LaTeX and generates a condensed outline."""
    llm = ChatGroq(model=model_name, api_key=GROQ_API_KEY, temperature=0)
    chunks = latex_parse.invoke(latex_file)
    
    summaries = []
    for i, chunk in enumerate(chunks):
        with st.status(f"Reading Section {i+1}...", expanded=False):
            res = llm.invoke(f"Summarize this LaTeX section for a presentation. Focus on key technical results: {chunk}")
            summaries.append(res.content)
    
    master_prompt = "Combine these summaries into a structured slide-by-slide outline (10-15 slides total):\n\n" + "\n".join(summaries)
    outline = llm.invoke(master_prompt)
    st.session_state.agent_thoughts += f"\n--- SUMMARIZER ---\nParsed {len(chunks)} chunks into a master outline."
    return outline.content

def run_creator(outline, model_name):
    """Agent 2: Converts the outline into the actual PPTX file."""
    llm = ChatGroq(model=model_name, api_key=GROQ_API_KEY, temperature=0)
    agent = initialize_agent([ppt_create], llm, agent=AgentType.STRUCTURED_CHAT_ZERO_SHOT_REACT_DESCRIPTION, verbose=True)
    
    f = io.StringIO()
    with contextlib.redirect_stdout(f):
        prompt = f"Using this outline, call ppt_create to build the PowerPoint: {outline}. MANDATORY: 'points' MUST be a list of strings."
        response = agent.run(prompt)
    st.session_state.agent_thoughts += f"\n--- CREATOR ---\n{f.getvalue()}"
    return response

def run_designer(ppt_path, model_name):
    """Agent 3: Beautifies the existing file."""
    llm = ChatGroq(model=model_name, api_key=GROQ_API_KEY, temperature=0)
    agent = initialize_agent([beautify_presentation], llm, agent=AgentType.STRUCTURED_CHAT_ZERO_SHOT_REACT_DESCRIPTION, verbose=True)
    
    f = io.StringIO()
    with contextlib.redirect_stdout(f):
        prompt = f"Decide the best theme for {ppt_path} and apply it using beautify_presentation."
        response = agent.run(prompt)
    st.session_state.agent_thoughts += f"\n--- DESIGNER ---\n{f.getvalue()}"
    return response

# --- APP UI ---
st.title("📄 Research Paper to PPTX")
st.caption("Upload a single `.tex` file to generate a themed presentation.")
model = st.selectbox("LLM Selection", ("llama-3.3-70b-versatile", "openai/gpt-oss-20b", "openai/gpt-oss-120b", "qwen/qwen3-32b"))
uploaded_file = st.file_uploader("Upload LaTeX Paper", type=["tex"])

if uploaded_file and st.session_state.step == "upload":
    with open("temp.tex", "wb") as f: f.write(uploaded_file.getbuffer())
    if st.button("Generate Outline"):
        with st.spinner("Summarizer is reading the paper..."):
            st.session_state.outline = run_summarizer("temp.tex", model)
            st.session_state.step = "approve"
            st.rerun()

if st.session_state.step == "approve":
    st.subheader("Step 1: Review Paper Outline")
    with st.expander("View Master Outline", expanded=True):
        st.markdown(st.session_state.outline)
    
    col1, col2 = st.columns(2)
    with col1:
        if st.button("Approve Outline & Create Slides"):
            with st.spinner("Creator is building PPTX..."):
                run_creator(st.session_state.outline, model)
                st.session_state.ppt_path = os.path.abspath("output_presentation.pptx")
                st.session_state.step = "design"
                st.rerun()
    with col2:
        if st.button("Restart"):
            st.session_state.clear()
            st.rerun()

if st.session_state.step == "design":
    st.subheader("Step 2: Finalize Design")
    st.success("Slides have been created. Click below to apply a theme and finish.")
    if st.button("✨ Beautify & Complete"):
        with st.spinner("Designer is styling..."):
            run_designer(st.session_state.ppt_path, model)
            st.session_state.step = "final"
            st.rerun()

if st.session_state.step == "final":
    st.balloons()
    with open(st.session_state.ppt_path, "rb") as f:
        st.download_button("Download Final Presentation", f, file_name="final_presentation.pptx")
    if st.button("Start New Project"):
        st.session_state.clear()
        st.rerun()
